"""
Action Executor: Comprehensive tool for executing external APIs, device control, system commands,
web automation, email, messaging, smart home, file operations, research, and more.
"""

import csv
import json
import logging
import os
import shutil
import smtplib
import subprocess
import sys
from datetime import datetime
from email.message import EmailMessage
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

import requests

from ..providers import ProviderLookupError, provider_registry

DEFAULT_SENSITIVE_ACTIONS = {
    "send_email",
    "send_whatsapp",
    "send_message",
    "send_telegram",
    "run_command",
    "run_sandboxed",
    "web_automation",
    "open_file",
    "open_directory",
    "create_presentation",
    "add_presentation_slide",
    "desktop_close_app",
    "desktop_click",
    "desktop_send_keys",
    "desktop_execute_schema",
}


class ActionExecutor:
    ACTION_SCOPES: Dict[str, set[str]] = {
        "open_file": {"files.read"},
        "read_file": {"files.read"},
        "search_file": {"files.read"},
        "create_file": {"files.write"},
        "write_file": {"files.write"},
        "delete_file": {"files.write"},
        "move_file": {"files.write"},
        "create_directory": {"files.write"},
        "open_directory": {"files.read"},
        "run_command": {"system"},
        "run_sandboxed": {"system"},
        "open_application": {"system"},
        "launch_app": {"system"},
        "close_app": {"system"},
        "clipboard_copy": {"system"},
        "clipboard_paste": {"system"},
        "clipboard": {"system"},
        "show_notification": {"system"},
        "notify": {"system"},
        "system_setting": {"system"},
        "get_system_info": {"system"},
        "send_email": {"messaging"},
        "send_whatsapp": {"messaging"},
        "send_message": {"messaging"},
        "send_telegram": {"messaging"},
        "web_search": {"web"},
        "web_scrape": {"web"},
        "web_automation": {"web"},
        "research_topic": {"web"},
        "wikipedia_search": {"web"},
        "control_device": {"iot"},
        "smart_home": {"iot"},
        "control_lights": {"iot"},
        "control_temperature": {"iot"},
        "store_memory": {"memory.write"},
        "search_memory": {"memory.read"},
        "create_code": {"files.write"},
        "analyze_code": {"files.read"},
        "make_note": {"files.write"},
        "read_notes": {"files.read"},
        "search_notes": {"files.read"},
        "create_docx": {"files.write"},
        "create_pdf": {"files.write"},
        "create_sheet": {"files.write"},
        "read_sheet": {"files.read"},
        "write_sheet": {"files.write"},
        "open_sheet": {"files.read"},
        "create_presentation": {"files.write"},
        "add_presentation_slide": {"files.write"},
        "open_presentation": {"files.read"},
        "calendar_event": {"productivity"},
        "desktop_open_app": {"desktop"},
        "desktop_close_app": {"desktop"},
        "desktop_type_text": {"desktop"},
        "desktop_click": {"desktop"},
        "desktop_send_keys": {"desktop"},
        "desktop_get_text": {"desktop"},
        "desktop_execute_schema": {"desktop"},
    }

    def __init__(
        self,
        permissions: Optional[Dict[str, bool]] = None,
        logger: Optional[logging.Logger] = None,
        consent_callback: Optional[
            Callable[[str, Dict[str, Any]], bool]
        ] = None,
        sensitive_actions: Optional[List[str]] = None,
        *,
        security_manager: Optional[Any] = None,
        rag_pipeline: Optional[Any] = None,
        web_agent: Optional[Any] = None,
        config_manager: Optional[Any] = None,
        allowed_scopes: Optional[List[str]] = None,
    ) -> None:
        self.permissions = permissions or {}
        self.logger = logger or logging.getLogger(__name__)
        self.consent_callback = consent_callback or (
            lambda action, params: True
        )
        self.sensitive_actions = set(
            sensitive_actions or DEFAULT_SENSITIVE_ACTIONS
        )
        self.notes_file = "mia_notes.md"
        self.security_manager = security_manager
        self.rag_pipeline = rag_pipeline
        self._config_manager = config_manager
        self.allowed_scopes = set(allowed_scopes or [])

        self.config = self._load_config()
        if self._config_manager and getattr(
            self._config_manager, "config", None
        ):
            try:
                cfg = self._config_manager.config
                if getattr(cfg, "documents", None):
                    self.config.setdefault("documents", {})
                    self.config["documents"].update(
                        {
                            "template_dir": getattr(
                                cfg.documents,
                                "template_dir",
                                self.config["documents"].get("template_dir"),
                            ),
                            "output_dir": getattr(
                                cfg.documents,
                                "output_dir",
                                self.config["documents"].get("output_dir"),
                            ),
                            "default_template": getattr(
                                cfg.documents,
                                "default_template",
                                self.config["documents"].get(
                                    "default_template"
                                ),
                            ),
                        }
                    )
                if getattr(cfg, "memory", None):
                    self.config.setdefault("memory", {})
                    self.config["memory"].update(
                        {
                            "persist_dir": getattr(
                                cfg.memory,
                                "vector_db_path",
                                self.config["memory"].get("persist_dir"),
                            ),
                            "vector_enabled": getattr(
                                cfg.memory,
                                "vector_enabled",
                                self.config["memory"].get(
                                    "vector_enabled", True
                                ),
                            ),
                            "graph_enabled": getattr(
                                cfg.memory,
                                "graph_enabled",
                                self.config["memory"].get(
                                    "graph_enabled", True
                                ),
                            ),
                            "long_term_enabled": getattr(
                                cfg.memory,
                                "long_term_enabled",
                                self.config["memory"].get(
                                    "long_term_enabled", True
                                ),
                            ),
                            "max_entries": getattr(
                                cfg.memory,
                                "max_memory_size",
                                self.config["memory"].get(
                                    "max_entries", 10_000
                                ),
                            ),
                            "max_results": getattr(
                                cfg.memory,
                                "max_results",
                                self.config["memory"].get("max_results", 5),
                            ),
                            "similarity_threshold": getattr(
                                cfg.memory,
                                "similarity_threshold",
                                self.config["memory"].get(
                                    "similarity_threshold", 0.7
                                ),
                            ),
                        }
                    )
            except Exception as exc:
                self.logger.debug(
                    "Failed to merge config manager settings into ActionExecutor: %s",
                    exc,
                )

        self._document_factory = self._load_provider_factory("documents")
        self._sandbox_factory = self._load_provider_factory("sandbox")
        self._telegram_factory = self._load_provider_factory(
            "messaging", "telegram"
        )
        self._memory_factory = self._load_provider_factory("memory", "manager")
        self._web_agent_factory = (
            self._load_provider_factory("web", "agent")
            if web_agent is None
            else None
        )
        self._desktop_factory = self._load_provider_factory("desktop", "automation")
        self._document_generator_instance: Optional[Any] = None
        self._sandbox_instance: Optional[Any] = None
        self._telegram_client: Optional[Any] = None
        self._memory_instance: Optional[Any] = None
        self._web_agent_instance: Optional[Any] = web_agent
        self._desktop_instance: Optional[Any] = None

    def _load_config(self) -> Dict[str, Any]:
        """Load configuration from config file or environment variables."""

        def _to_int(env_var: str, default: int) -> int:
            value = os.getenv(env_var)
            if value is None:
                return default
            try:
                return int(value)
            except ValueError:
                self.logger.warning(
                    "Invalid integer for %s: %s", env_var, value
                )
                return default

        def _env(*names: str, default: Optional[str] = None) -> Optional[str]:
            for name in names:
                value = os.getenv(name)
                if value not in (None, ""):
                    return value
            return default

        def _to_int_env(
            *names: str, default: Optional[int] = None
        ) -> Optional[int]:
            raw = _env(*names)
            if raw in (None, ""):
                return default
            try:
                return int(str(raw))
            except ValueError:
                self.logger.warning(
                    "Invalid integer value for %s: %s", names[0], raw
                )
                return default

        def _to_bool_env(*names: str, default: bool = False) -> bool:
            raw = _env(*names)
            if raw is None:
                return default
            return str(raw).strip().lower() in {"1", "true", "yes", "on"}

        def _to_float_env(
            *names: str, default: Optional[float] = None
        ) -> Optional[float]:
            raw = _env(*names)
            if raw in (None, ""):
                return default
            try:
                return float(str(raw))
            except ValueError:
                self.logger.warning(
                    "Invalid float value for %s: %s", names[0], raw
                )
                return default

        config = {
            "email": {
                "smtp_server": os.getenv("SMTP_SERVER", "smtp.gmail.com"),
                "smtp_port": int(os.getenv("SMTP_PORT", "587")),
                "username": os.getenv("EMAIL_USERNAME", ""),
                "password": os.getenv("EMAIL_PASSWORD", ""),
            },
            "research": {
                "google_api_key": os.getenv("GOOGLE_API_KEY", ""),
                "google_cse_id": os.getenv("GOOGLE_CSE_ID", ""),
                "default_search_engine": "duckduckgo",
            },
            "smart_home": {
                "home_assistant_url": os.getenv("HOME_ASSISTANT_URL", ""),
                "home_assistant_token": os.getenv("HOME_ASSISTANT_TOKEN", ""),
            },
            "telegram": {
                "enabled": _to_bool_env(
                    "MIA_TELEGRAM_ENABLED", "TELEGRAM_ENABLED"
                ),
                "api_id": _to_int_env(
                    "MIA_TELEGRAM_API_ID", "TELEGRAM_API_ID"
                ),
                "api_hash": _env(
                    "MIA_TELEGRAM_API_HASH", "TELEGRAM_API_HASH", default=""
                ),
                "bot_token": _env(
                    "MIA_TELEGRAM_BOT_TOKEN", "TELEGRAM_BOT_TOKEN", default=""
                ),
                "phone_number": _env("MIA_TELEGRAM_PHONE", "TELEGRAM_PHONE"),
                "session_dir": _env(
                    "MIA_TELEGRAM_SESSION_DIR",
                    "TELEGRAM_SESSION_DIR",
                    default="sessions",
                ),
                "session_name": _env(
                    "MIA_TELEGRAM_SESSION_NAME",
                    "TELEGRAM_SESSION_NAME",
                    default="mia_telegram",
                ),
                "default_peer": _env(
                    "MIA_TELEGRAM_DEFAULT_PEER",
                    "TELEGRAM_DEFAULT_PEER",
                    default="",
                ),
                "parse_mode": _env(
                    "MIA_TELEGRAM_PARSE_MODE",
                    "TELEGRAM_PARSE_MODE",
                    default="markdown",
                ),
                "request_timeout": _to_int_env(
                    "MIA_TELEGRAM_REQUEST_TIMEOUT",
                    "TELEGRAM_REQUEST_TIMEOUT",
                    default=30,
                ),
            },
            "whatsapp": {"phone_number": os.getenv("WHATSAPP_PHONE", "")},
            "sandbox": {
                "work_dir": os.getenv("MIA_SANDBOX_WORKDIR", "sandbox_runs"),
                "log_dir": os.getenv("MIA_SANDBOX_LOGDIR", "logs/sandbox"),
                "memory_mb": _to_int("MIA_SANDBOX_MEMORY_MB", 256),
                "timeout_ms": _to_int("MIA_SANDBOX_TIMEOUT_MS", 10_000),
                "fuel": os.getenv("MIA_SANDBOX_FUEL"),
            },
            "documents": {
                "template_dir": os.getenv(
                    "MIA_DOC_TEMPLATE_DIR", "templates/documents"
                ),
                "output_dir": os.getenv(
                    "MIA_DOC_OUTPUT_DIR", "output/documents"
                ),
                "default_template": os.getenv(
                    "MIA_DOC_DEFAULT_TEMPLATE", "proposal"
                ),
            },
            "memory": {
                "persist_dir": _env("MIA_MEMORY_PATH", default="memory"),
                "vector_enabled": _to_bool_env(
                    "MIA_MEMORY_VECTOR_ENABLED", default=True
                ),
                "graph_enabled": _to_bool_env(
                    "MIA_MEMORY_GRAPH_ENABLED", default=True
                ),
                "long_term_enabled": _to_bool_env(
                    "MIA_MEMORY_LONG_TERM_ENABLED", default=True
                ),
                "max_entries": _to_int_env(
                    "MIA_MEMORY_MAX_ENTRIES", default=10_000
                )
                or 10_000,
                "max_results": _to_int_env("MIA_MEMORY_MAX_RESULTS", default=5)
                or 5,
                "similarity_threshold": _to_float_env(
                    "MIA_MEMORY_SIMILARITY", default=0.7
                )
                or 0.7,
            },
        }
        return config

    def _load_provider_factory(self, domain: str, name: Optional[str] = None):
        try:
            return provider_registry.get_factory(domain, name)
        except ProviderLookupError:
            self.logger.debug("Provider not registered: %s", domain)
            return None
        except Exception as exc:  # pragma: no cover
            self.logger.warning("Failed to load provider %s: %s", domain, exc)
            return None

    def _get_document_generator(self):
        if self._document_generator_instance is False:
            return None
        if self._document_generator_instance is not None:
            return self._document_generator_instance
        if not self._document_factory:
            self._document_generator_instance = False
            return None
        try:
            self._document_generator_instance = self._document_factory(
                template_dir=self.config["documents"]["template_dir"],
                output_dir=self.config["documents"]["output_dir"],
                logger_instance=self.logger,
            )
        except Exception as exc:  # pragma: no cover
            self.logger.warning(
                "Failed to instantiate document generator: %s", exc
            )
            self._document_generator_instance = False
        return self._document_generator_instance or None

    def _get_sandbox(self):
        if self._sandbox_instance is False:
            return None
        if self._sandbox_instance is not None:
            return self._sandbox_instance
        if not self._sandbox_factory:
            self._sandbox_instance = False
            return None
        try:
            from ..sandbox.wasi_runner import SandboxLimits  # type: ignore

            fuel_value = self.config["sandbox"].get("fuel")
            try:
                fuel_int = (
                    int(fuel_value) if fuel_value not in (None, "") else None
                )
            except ValueError:
                self.logger.warning(
                    "Invalid sandbox fuel value: %s", fuel_value
                )
                fuel_int = None
            limits = SandboxLimits(
                max_memory_mb=self.config["sandbox"].get("memory_mb", 256),
                timeout_ms=self.config["sandbox"].get("timeout_ms", 10_000),
                fuel=fuel_int,
            )
            self._sandbox_instance = self._sandbox_factory(
                limits=limits,
                work_dir=self.config["sandbox"].get("work_dir"),
                log_dir=self.config["sandbox"].get("log_dir"),
            )
        except Exception as exc:  # pragma: no cover
            self.logger.warning("Failed to instantiate sandbox: %s", exc)
            self._sandbox_instance = False
        return self._sandbox_instance or None

    def _get_memory_manager(self):
        if self._memory_instance is False:
            return None
        if self._memory_instance is not None:
            return self._memory_instance
        if not self._memory_factory:
            self._memory_instance = False
            return None
        try:
            memory_cfg = self.config.get("memory", {})
            self._memory_instance = self._memory_factory(
                persist_directory=memory_cfg.get("persist_dir", "memory"),
                vector_enabled=memory_cfg.get("vector_enabled", True),
                graph_enabled=memory_cfg.get("graph_enabled", True),
                long_term_enabled=memory_cfg.get("long_term_enabled", True),
                max_long_term_entries=memory_cfg.get("max_entries", 10_000),
                max_results=memory_cfg.get("max_results", 5),
                similarity_threshold=memory_cfg.get(
                    "similarity_threshold", 0.7
                ),
                logger_instance=self.logger,
            )
        except Exception as exc:  # pragma: no cover
            self.logger.warning(
                "Failed to instantiate memory manager: %s", exc
            )
            self._memory_instance = False
        return self._memory_instance or None

    def _get_web_agent(self):
        if self._web_agent_instance is False:
            return None
        if self._web_agent_instance is not None:
            return self._web_agent_instance
        if not self._web_agent_factory:
            self._web_agent_instance = False
            return None
        try:
            self._web_agent_instance = self._web_agent_factory()
        except Exception as exc:  # pragma: no cover - optional deps
            self.logger.warning("Failed to instantiate web agent: %s", exc)
            self._web_agent_instance = False
        return self._web_agent_instance or None

    def _get_desktop_automation(self):
        if self._desktop_instance is False:
            return None
        if self._desktop_instance is not None:
            return self._desktop_instance
        if not self._desktop_factory:
            self._desktop_instance = False
            return None
        try:
            self._desktop_instance = self._desktop_factory(logger_instance=self.logger)
        except Exception as exc:  # pragma: no cover - optional deps
            self.logger.warning("Failed to instantiate desktop automation: %s", exc)
            self._desktop_instance = False
        return self._desktop_instance or None

    def _is_scope_allowed(self, action: str) -> bool:
        required = self.ACTION_SCOPES.get(action)
        if not required:
            return True

        if self.allowed_scopes:
            if "*" in self.allowed_scopes or "all" in self.allowed_scopes:
                return True
            if required & self.allowed_scopes:
                return True
            return False

        if self.security_manager:
            try:
                if hasattr(self.security_manager, "has_scopes"):
                    return bool(self.security_manager.has_scopes(required))  # type: ignore[attr-defined]
                if hasattr(self.security_manager, "has_scope"):
                    return any(
                        bool(self.security_manager.has_scope(scope))  # type: ignore[attr-defined]
                        for scope in required
                    )
            except Exception as exc:  # pragma: no cover - defensive
                self.logger.warning(
                    "Security scope check failed for %s: %s", action, exc
                )
                return False

        return True

    def _get_telegram_client(self):
        if self._telegram_client is False:
            return None
        if self._telegram_client is not None:
            return self._telegram_client
        if not self._telegram_factory:
            self._telegram_client = False
            return None
        try:
            telegram_cfg = self.config.get("telegram", {})
            self._telegram_client = self._telegram_factory(
                api_id=telegram_cfg.get("api_id"),
                api_hash=telegram_cfg.get("api_hash"),
                bot_token=telegram_cfg.get("bot_token"),
                phone_number=telegram_cfg.get("phone_number"),
                session_dir=telegram_cfg.get("session_dir"),
                session_name=telegram_cfg.get("session_name"),
                default_peer=telegram_cfg.get("default_peer"),
                parse_mode=telegram_cfg.get("parse_mode"),
                request_timeout=telegram_cfg.get("request_timeout", 30),
                enabled=telegram_cfg.get("enabled", False),
                logger_instance=self.logger,
            )
        except Exception as exc:  # pragma: no cover
            self.logger.warning(
                "Failed to instantiate Telegram messenger: %s", exc
            )
            self._telegram_client = False
        return self._telegram_client or None

    def execute(self, action: str, params: Optional[Dict[str, Any]]):
        """Execute an action with given parameters."""
        if not action:
            return "No action provided."

        params = params or {}
        self.logger.info(f"Executing action: {action}")

        consent_prompted = False

        if not self._is_scope_allowed(action):
            self.logger.warning(
                "Action %s blocked due to missing scopes", action
            )
            return f"Action '{action}' is not permitted in the current security scope."

        if self.security_manager:
            try:
                if hasattr(self.security_manager, "log_action"):
                    self.security_manager.log_action(action, params)
                elif hasattr(self.security_manager, "record_action"):
                    self.security_manager.record_action(action, params)
            except Exception as exc:  # pragma: no cover - defensive
                self.logger.debug(
                    "Security action logging failed for %s: %s", action, exc
                )

        # Check permissions (allow all actions by default unless specifically restricted)
        if (
            self.permissions
            and action in self.permissions
            and not self.permissions[action]
        ):
            consent_prompted = True
            if not self.consent_callback(action, params):
                self.logger.warning(f"Permission denied for action: {action}")
                return f"Permission denied for action: {action}"

        if action in self.sensitive_actions and not consent_prompted:
            if not self.consent_callback(action, params):
                self.logger.info(
                    "Action cancelled after consent prompt: %s", action
                )
                return f"Action '{action}' cancelled by user consent."

        try:
            # Use dispatch pattern for better maintainability
            action_dispatch = {
                # File operations
                "open_file": lambda: self.open_file(params.get("path") or ""),
                "create_file": lambda: self.create_file(
                    params.get("path") or "", params.get("content", "")
                ),
                "read_file": lambda: self.read_file(params.get("path") or ""),
                "write_file": lambda: self.write_file(
                    params.get("path") or "", params.get("content", "")
                ),
                "move_file": lambda: self.move_file(
                    params.get("src"), params.get("dst")
                ),
                "delete_file": lambda: self.delete_file(
                    params.get("path") or ""
                ),
                "search_file": lambda: self.search_file(
                    params.get("name") or "", params.get("directory", ".")
                ),
                "open_directory": lambda: self.open_directory(
                    params.get("path") or ""
                ),
                "create_directory": lambda: self.create_directory(
                    params.get("path") or ""
                ),
                # Code generation
                "create_code": lambda: self.create_code(
                    params.get("language") or "",
                    params.get("description") or "",
                    params.get("filename"),
                ),
                "analyze_code": lambda: self.analyze_code(
                    params.get("path") or ""
                ),
                # Notes and documentation
                "make_note": lambda: self.make_note(
                    params.get("content") or "", params.get("title")
                ),
                "read_notes": lambda: self.read_notes(),
                "search_notes": lambda: self.search_notes(
                    params.get("query") or ""
                ),
                "create_docx": lambda: self.generate_docx(params),
                "create_pdf": lambda: self.generate_pdf(params),
                "store_memory": lambda: self.store_memory(params),
                "search_memory": lambda: self.search_memory(params),
                "link_memory_nodes": lambda: self.link_memory_nodes(params),
                # Spreadsheet operations
                "create_sheet": lambda: self.create_sheet(
                    params.get("filename") or "", params.get("data")
                ),
                "read_sheet": lambda: self.read_sheet(
                    params.get("filename") or ""
                ),
                "write_sheet": lambda: self.write_sheet(
                    params.get("filename") or "", params.get("data") or []
                ),
                "open_sheet": lambda: self.open_file(
                    params.get("filename") or ""
                ),
                # Presentation (PowerPoint)
                "create_presentation": lambda: self.create_presentation(
                    params.get("filename") or "",
                    params.get("title"),
                    params.get("content"),
                ),
                "add_presentation_slide": lambda: self.add_presentation_slide(
                    params.get("filename") or "",
                    params.get("title"),
                    params.get("bullets") or [],
                ),
                "open_presentation": lambda: self.open_file(
                    params.get("filename") or ""
                ),
                # Research and web operations
                "web_search": lambda: self.web_search(
                    params.get("query") or ""
                ),
                "web_scrape": lambda: self.web_scrape(params.get("url") or ""),
                "research_topic": lambda: self.research_topic(
                    params.get("topic") or ""
                ),
                "wikipedia_search": lambda: self.wikipedia_search(
                    params.get("query") or ""
                ),
                # Smart home integration
                "control_device": lambda: self._handle_control_device(params),
                # System integration
                "clipboard_copy": lambda: self.clipboard_copy(
                    params.get("text") or ""
                ),
                "clipboard_paste": lambda: self.clipboard_paste(),
                "show_notification": lambda: self.show_notification(
                    params.get("title") or "", params.get("message") or ""
                ),
                "open_application": lambda: self.open_application(
                    params.get("app_name") or ""
                ),
                "get_system_info": lambda: self.get_system_info(),
                # Application control
                "launch_app": lambda: self.launch_app(params.get("app") or ""),
                "close_app": lambda: self.close_app(params.get("app") or ""),
                "clipboard": lambda: self.clipboard_action(params),
                "notify": lambda: self.notify(params.get("message")),
                "system_setting": lambda: self.system_setting(params),
                "run_command": lambda: self.run_command(params.get("command")),
                "run_sandboxed": lambda: self.run_sandboxed(params),
                "web_automation": lambda: self.web_automation(params),
                # Communication
                "send_email": lambda: self.send_email(params),
                "send_whatsapp": lambda: self.send_whatsapp(params),
                "send_message": lambda: self.send_message(params),
                "send_telegram": lambda: self._send_telegram(
                    params.get("to", ""), params.get("message", "")
                ),
                # Calendar and scheduling
                "calendar_event": lambda: self.calendar_event(params),
                # Smart home
                "smart_home": lambda: self.smart_home(params),
                "control_lights": lambda: self.control_lights(params),
                "control_temperature": lambda: self.control_temperature(
                    params
                ),
                # Desktop automation
                "desktop_open_app": lambda: self.desktop_open_app(params),
                "desktop_close_app": lambda: self.desktop_close_app(params),
                "desktop_type_text": lambda: self.desktop_type_text(params),
                "desktop_click": lambda: self.desktop_click(params),
                "desktop_send_keys": lambda: self.desktop_send_keys(params),
                "desktop_get_text": lambda: self.desktop_get_text(params),
                "desktop_execute_schema": lambda: self.desktop_execute_schema(params),
            }

            # Execute the action if it exists
            if action in action_dispatch:
                return action_dispatch[action]()
            else:
                self.logger.error(f"Unknown action: {action}")
                return f"Unknown action: {action}"

        except Exception as e:
            self.logger.error(f"Error executing {action}: {e}")
            return f"Error: {e}"

    def _handle_control_device(self, params):
        """Handle device control with parameter filtering."""
        device_type = params.get("device_type")
        device_action = params.get("action")
        # Remove these from params to avoid duplicate keyword arguments
        filtered_params = {
            k: v
            for k, v in params.items()
            if k not in ["device_type", "action"]
        }
        return self.control_device(
            device_type, device_action, **filtered_params
        )

    # Enhanced File Operations
    def create_file(self, path: str, content: str = "") -> str:
        """Create a new file with specified content."""
        if not path:
            return "No file path provided."
        try:
            Path(path).parent.mkdir(parents=True, exist_ok=True)
            with open(path, "w", encoding="utf-8") as f:
                f.write(content)
            return f"Created file: {path}"
        except (OSError, IOError) as e:
            return f"Error creating file: {e}"
        except UnicodeEncodeError as e:
            return f"Error encoding content: {e}"
        except Exception as e:
            return f"Unexpected error creating file: {e}"

    def read_file(self, path: str) -> str:
        """Read content from a file."""
        if not path:
            return "No file path provided."
        try:
            with open(path, "r", encoding="utf-8") as f:
                content = f.read()
            return f"Content of {path}:\n{content}"
        except FileNotFoundError:
            return f"File not found: {path}"
        except PermissionError:
            return f"Permission denied reading file: {path}"
        except (OSError, IOError) as e:
            return f"Error reading file: {e}"
        except UnicodeDecodeError as e:
            return f"Error decoding file content: {e}"
        except Exception as e:
            return f"Unexpected error reading file: {e}"

    def write_file(self, path: str, content: str) -> str:
        """Write content to a file."""
        if not path:
            return "No file path provided."
        try:
            with open(path, "w", encoding="utf-8") as f:
                f.write(content)
            return f"Written to file: {path}"
        except PermissionError:
            return f"Permission denied writing to file: {path}"
        except (OSError, IOError) as e:
            return f"Error writing file: {e}"
        except UnicodeEncodeError as e:
            return f"Error encoding content: {e}"
        except Exception as e:
            return f"Unexpected error writing file: {e}"

    def open_directory(self, path: str) -> str:
        """Open a directory in the file explorer."""
        if not path:
            return "No directory path provided."
        try:
            if os.name == "nt":  # Windows
                os.startfile(path)
            elif os.name == "posix":  # macOS and Linux
                subprocess.run(
                    ["open", path]
                    if sys.platform == "darwin"
                    else ["xdg-open", path]
                )
            return f"Opened directory: {path}"
        except FileNotFoundError:
            return f"Directory not found: {path}"
        except PermissionError:
            return f"Permission denied opening directory: {path}"
        except (OSError, IOError) as e:
            return f"Error opening directory: {e}"
        except subprocess.SubprocessError as e:
            return f"Error running file manager: {e}"
        except Exception as e:
            return f"Unexpected error opening directory: {e}"

    def create_directory(self, path: str) -> str:
        """Create a new directory."""
        if not path:
            return "No directory path provided."
        try:
            Path(path).mkdir(parents=True, exist_ok=True)
            return f"Created directory: {path}"
        except Exception as e:
            return f"Error creating directory: {e}"

    # Document generation
    def generate_docx(self, params: Dict[str, Any]) -> str:
        generator = self._get_document_generator()
        if not generator:
            return "Document generator not available. Install python-docx and enable the provider."

        context_input = params.get("context")
        context = (
            dict(context_input) if isinstance(context_input, dict) else {}
        )
        summary = (
            params.get("summary")
            or params.get("content")
            or params.get("body")
        )
        title = params.get("title") or params.get("name")
        if summary:
            context.setdefault("summary", summary)
        if title:
            context.setdefault("title", title)

        template = params.get("template") or self.config["documents"].get(
            "default_template", "proposal"
        )
        output_path = params.get("output_path")
        return generator.create_docx(template, context, output_path)

    def generate_pdf(self, params: Dict[str, Any]) -> str:
        generator = self._get_document_generator()
        if not generator:
            return "Document generator not available. Install reportlab and enable the provider."

        context_input = params.get("context")
        context = (
            dict(context_input) if isinstance(context_input, dict) else {}
        )
        summary = (
            params.get("summary")
            or params.get("content")
            or params.get("body")
        )
        title = params.get("title") or params.get("name")
        if summary:
            context.setdefault("summary", summary)
        if title:
            context.setdefault("title", title)

        template = params.get("template") or self.config["documents"].get(
            "default_template", "proposal"
        )
        output_path = params.get("output_path")
        return generator.create_pdf(template, context, output_path)

    # Memory operations ----------------------------------------------------

    def store_memory(self, params: Dict[str, Any]) -> str:
        manager = self._get_memory_manager()
        if not manager:
            return "Memory manager not available. Enable memory support and install optional dependencies."

        text = (
            params.get("text")
            or params.get("content")
            or params.get("summary")
        )
        if not text:
            return "No memory content provided."

        embedding = params.get("embedding")
        if isinstance(embedding, str):
            try:
                cleaned = embedding.replace("[", "").replace("]", "")
                embedding = [
                    float(item.strip())
                    for item in cleaned.split(",")
                    if item.strip()
                ]
            except ValueError:
                self.logger.warning(
                    "Invalid embedding string provided; ignoring embedding"
                )
                embedding = None
        metadata = (
            params.get("metadata")
            if isinstance(params.get("metadata"), dict)
            else {}
        )
        relations_input = params.get("relations")
        relation_entries = None
        if relations_input is not None:
            relation_entries = []
            items = (
                relations_input
                if isinstance(relations_input, (list, tuple))
                else [relations_input]
            )
            for item in items:
                if isinstance(item, dict):
                    source = item.get("source")
                    target = item.get("target")
                    if not source or not target:
                        continue
                    relation_entries.append(
                        (
                            source,
                            item.get("relation")
                            or item.get("predicate")
                            or "related_to",
                            target,
                            item.get("metadata"),
                        )
                    )
                elif isinstance(item, (list, tuple)):
                    if len(item) >= 3:
                        relation_entries.append(tuple(item))
            if not relation_entries:
                relation_entries = None

        # Store in RAG pipeline if available
        rag_result = None
        if self.rag_pipeline and hasattr(self.rag_pipeline, "remember"):
            try:
                enriched_metadata = dict(metadata) if metadata else {}
                enriched_metadata.setdefault("source", "action_executor")
                enriched_metadata.setdefault(
                    "timestamp", datetime.now().isoformat()
                )
                rag_result = self.rag_pipeline.remember(
                    text, metadata=enriched_metadata
                )
                self.logger.debug(
                    "Stored memory in RAG pipeline: %s", rag_result
                )
            except Exception as exc:
                self.logger.debug("Failed to store in RAG pipeline: %s", exc)

        try:
            record = manager.store_memory(
                text,
                embedding=embedding,
                metadata=metadata,
                relations=relation_entries,
            )
            reference = record.get("vector_id") or record.get("long_term_id")
            result_msg = (
                f"Memory stored successfully (ref={reference or 'n/a'})."
            )
            if rag_result:
                result_msg += f" Also stored in RAG pipeline (id={rag_result.get('id', 'n/a')})."
            return result_msg
        except Exception as exc:
            self.logger.error("Failed to store memory: %s", exc)
            return f"Error storing memory: {exc}"

    def search_memory(self, params: Dict[str, Any]) -> str:
        # Try RAG pipeline first if available
        if self.rag_pipeline and hasattr(self.rag_pipeline, "query"):
            query = (
                params.get("query")
                or params.get("text")
                or params.get("keyword")
            )
            if query:
                try:
                    top_k_raw = params.get("top_k") or params.get("limit", 4)
                    top_k = (
                        int(top_k_raw)
                        if isinstance(top_k_raw, (int, str))
                        and str(top_k_raw).isdigit()
                        else 4
                    )

                    rag_results = self.rag_pipeline.query(query, top_k=top_k)
                    if rag_results:
                        lines: List[str] = []
                        for idx, chunk in enumerate(rag_results, start=1):
                            preview = chunk.text
                            if len(preview) > 160:
                                preview = preview[:157] + "..."
                            score = chunk.score
                            metadata_info = []
                            if chunk.metadata.get("source"):
                                metadata_info.append(
                                    f"source: {chunk.metadata['source']}"
                                )
                            if chunk.metadata.get("timestamp"):
                                metadata_info.append(
                                    f"time: {chunk.metadata['timestamp']}"
                                )

                            segment = f"{idx}. {preview}"
                            if isinstance(score, (int, float)):
                                segment += f" (score={score:.3f})"
                            if metadata_info:
                                segment += f" [{', '.join(metadata_info)}]"
                            lines.append(segment)
                        return (
                            f"RAG Memory results for '{query}':\n"
                            + "\n".join(lines)
                        )
                except Exception as exc:
                    self.logger.debug(
                        "RAG pipeline search failed, falling back to memory manager: %s",
                        exc,
                    )

        # Fallback to memory manager
        manager = self._get_memory_manager()
        if not manager:
            return "Memory manager not available."

        query = (
            params.get("query") or params.get("text") or params.get("keyword")
        )
        embedding = params.get("embedding")
        if isinstance(embedding, str):
            cleaned = embedding.replace("[", "").replace("]", "")
            try:
                embedding = [
                    float(item.strip())
                    for item in cleaned.split(",")
                    if item.strip()
                ]
            except ValueError:
                self.logger.warning(
                    "Invalid embedding string provided for search; ignoring embedding"
                )
                embedding = None
        top_k_raw = params.get("top_k") or params.get("limit")
        top_k = None
        if top_k_raw is not None:
            try:
                top_k = int(top_k_raw)
            except (TypeError, ValueError):
                self.logger.warning(
                    "Invalid top_k value provided: %s", top_k_raw
                )
        try:
            results = manager.search_memory(
                embedding=embedding,
                query=query,
                top_k=top_k,
            )
        except Exception as exc:
            self.logger.error("Failed to search memory: %s", exc)
            return f"Error searching memory: {exc}"

        if not results:
            return "No memory results found."

        lines: List[str] = []
        for idx, item in enumerate(results, start=1):
            preview = str(item.get("text", ""))
            if len(preview) > 160:
                preview = preview[:157] + "..."
            source = item.get("source", "unknown")
            score = item.get("score")
            segment = f"{idx}. [{source}] {preview}"
            if isinstance(score, (int, float)):
                segment += f" (score={score:.3f})"
            lines.append(segment)
        return "Memory results:\n" + "\n".join(lines)

    def link_memory_nodes(self, params: Dict[str, Any]) -> str:
        manager = self._get_memory_manager()
        if not manager:
            return "Memory manager not available."

        source = params.get("source") or params.get("from")
        target = params.get("target") or params.get("to")
        relation = (
            params.get("relation") or params.get("predicate") or "related_to"
        )
        metadata = (
            params.get("metadata")
            if isinstance(params.get("metadata"), dict)
            else None
        )

        if not source or not target:
            return "Source and target are required to link memory nodes."

        try:
            manager.add_relation(
                str(source), str(relation), str(target), metadata=metadata
            )
            return f"Linked memory nodes: {source} -> {target} ({relation})."
        except Exception as exc:
            self.logger.error("Failed to link memory nodes: %s", exc)
            return f"Error linking memory nodes: {exc}"

    # Code Generation and Analysis
    def create_code(
        self, language: str, description: str, filename: Optional[str] = None
    ) -> str:
        """Generate code based on description and language."""
        if not language or not description:
            return "Language and description required."

        # Basic code templates
        templates = {
            "python": f"""# {description}
# Generated by M.I.A

def main():
    print("Hello, World!")
    # Add your code here

if __name__ == "__main__":
    main()
""",
            "javascript": f"""// {description}
// Generated by M.I.A

function main() {{
    console.log("Hello, World!");
    // Add your code here
}}
main();
""",
            "java": f"""// {description}
// Generated by M.I.A

public class Main {{
    public static void main(String[] args) {{
        System.out.println("Hello, World!");
        // Add your code here
    }}
}}
""",
            "html": f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{description}</title>
</head>
<body>
    <h1>{description}</h1>
    <!-- Add your HTML here -->
</body>
</html>
""",
            "css": f"""/* {description} */
/* Generated by M.I.A */

body {{
    font-family: Arial, sans-serif;
    margin: 0;
    padding: 20px;
}}

/* Add your CSS here */
""",
        }

        code = templates.get(
            language.lower(),
            f"// {description}\n// Code template not available for {language}",
        )

        if filename:
            try:
                with open(filename, "w", encoding="utf-8") as f:
                    f.write(code)
                return f"Code created and saved to {filename}"
            except Exception as e:
                return f"Error saving code: {e}"

        return f"Generated {language} code:\n{code}"

    def analyze_code(self, path: str) -> str:
        """Analyze code file and provide insights."""
        if not path:
            return "No file path provided."
        try:
            with open(path, "r", encoding="utf-8") as f:
                code = f.read()

            lines = code.split("\n")
            analysis = {
                "total_lines": len(lines),
                "non_empty_lines": len(
                    [line for line in lines if line.strip()]
                ),
                "comment_lines": len(
                    [
                        line
                        for line in lines
                        if line.strip().startswith("#")
                        or line.strip().startswith("//")
                    ]
                ),
                "file_size": os.path.getsize(path),
                "language": self._detect_language(path),
            }

            return (
                f"Code analysis for {path}:\n{json.dumps(analysis, indent=2)}"
            )
        except Exception as e:
            return f"Error analyzing code: {e}"

    def _detect_language(self, path: str) -> str:
        """Detect programming language from file extension."""
        extension = Path(path).suffix.lower()
        lang_map = {
            ".py": "Python",
            ".js": "JavaScript",
            ".java": "Java",
            ".cpp": "C++",
            ".c": "C",
            ".html": "HTML",
            ".css": "CSS",
            ".php": "PHP",
            ".rb": "Ruby",
            ".go": "Go",
            ".rs": "Rust",
        }
        return lang_map.get(extension, "Unknown")

    def move_file(self, src, dst):
        if not src or not dst:
            return "Source or destination missing."
        shutil.move(src, dst)
        return f"Moved {src} to {dst}"

    def delete_file(self, path):
        if not path:
            return "No file path provided."
        os.remove(path)
        return f"Deleted file: {path}"

    def search_file(self, name, directory="."):
        matches = []
        for root, dirs, files in os.walk(directory):
            if name in files:
                matches.append(os.path.join(root, name))
        return matches or f"No file named {name} found."

    # Notes and Documentation
    def make_note(self, content: str, title: Optional[str] = None) -> str:
        """Create or append to notes file."""
        if not content:
            return "No content provided for note."

        try:
            timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            note_entry = (
                f"\n## {title or 'Note'} - {timestamp}\n\n{content}\n\n---\n"
            )

            with open(self.notes_file, "a", encoding="utf-8") as f:
                f.write(note_entry)

            return f"Note saved to {self.notes_file}"
        except Exception as e:
            return f"Error saving note: {e}"

    def read_notes(self) -> str:
        """Read all notes from the notes file."""
        try:
            if not os.path.exists(self.notes_file):
                return "No notes file found."

            with open(self.notes_file, "r", encoding="utf-8") as f:
                notes = f.read()

            return f"Notes from {self.notes_file}:\n{notes}"
        except Exception as e:
            return f"Error reading notes: {e}"

    def search_notes(self, query: str) -> str:
        """Search for specific content in notes."""
        if not query:
            return "No search query provided."

        try:
            if not os.path.exists(self.notes_file):
                return "No notes file found."

            with open(self.notes_file, "r", encoding="utf-8") as f:
                notes = f.read()

            matching_lines = []
            for i, line in enumerate(notes.split("\n"), 1):
                if query.lower() in line.lower():
                    matching_lines.append(f"Line {i}: {line}")

            if matching_lines:
                return (
                    f"Found {len(matching_lines)} matches for '{query}':\n"
                    + "\n".join(matching_lines)
                )
            else:
                return f"No matches found for '{query}'"
        except Exception as e:
            return f"Error searching notes: {e}"

    # Spreadsheet Operations
    def create_sheet(
        self, filename: str, data: Optional[List[List[str]]] = None
    ) -> str:
        """Create a new spreadsheet file."""
        if not filename:
            return "No filename provided."

        try:
            if filename.endswith(".xlsx"):
                return self._create_excel_sheet(filename, data)
            elif filename.endswith(".csv"):
                return self._create_csv_sheet(filename, data)
            else:
                return "Unsupported file format. Use .xlsx or .csv"
        except Exception as e:
            return f"Error creating sheet: {e}"

    def _create_excel_sheet(
        self, filename: str, data: Optional[List[List[str]]] = None
    ) -> str:
        """Create Excel spreadsheet."""
        try:
            from openpyxl import Workbook  # type: ignore
        except ImportError:
            return "openpyxl not installed. Run: pip install openpyxl"

        wb = Workbook()  # type: ignore
        ws = wb.active  # type: ignore
        if ws is None:
            return "Failed to create worksheet"
        ws.title = "Sheet1"  # type: ignore

        if data:
            for row in data:
                ws.append(row)  # type: ignore
        else:
            ws.append(["Column1", "Column2", "Column3"])  # type: ignore
            ws.append(["Sample", "Data", "Here"])  # type: ignore

        wb.save(filename)
        return f"Excel sheet created: {filename}"

    def _create_csv_sheet(
        self, filename: str, data: Optional[List[List[str]]] = None
    ) -> str:
        """Create CSV file."""
        with open(filename, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            if data:
                writer.writerows(data)
            else:
                writer.writerow(["Column1", "Column2", "Column3"])
                writer.writerow(["Sample", "Data", "Here"])

        return f"CSV sheet created: {filename}"

    def read_sheet(self, filename: str) -> str:
        """Read data from spreadsheet file."""
        if not filename:
            return "No filename provided."

        try:
            if filename.endswith(".xlsx"):
                return self._read_excel_sheet(filename)
            elif filename.endswith(".csv"):
                return self._read_csv_sheet(filename)
            else:
                return "Unsupported file format. Use .xlsx or .csv"
        except Exception as e:
            return f"Error reading sheet: {e}"

    def _read_excel_sheet(self, filename: str) -> str:
        """Read Excel spreadsheet."""
        try:
            from openpyxl import load_workbook  # type: ignore
        except ImportError:
            return "openpyxl not installed. Run: pip install openpyxl"

        wb = load_workbook(filename)  # type: ignore
        ws = wb.active  # type: ignore
        if ws is None:
            return "Failed to load worksheet"

        data = []
        for row in ws.iter_rows(values_only=True):  # type: ignore
            data.append(row)

        return f"Excel sheet data from {filename}:\n{json.dumps(data, indent=2, default=str)}"

    def _read_csv_sheet(self, filename: str) -> str:
        """Read CSV file."""
        data = []
        with open(filename, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                data.append(row)

        return f"CSV sheet data from {filename}:\n{json.dumps(data, indent=2)}"

    def write_sheet(self, filename: str, data: List[List[str]]) -> str:
        """Write data to spreadsheet file."""
        if not filename or not data:
            return "Filename and data required."

        try:
            if filename.endswith(".xlsx"):
                return self._write_excel_sheet(filename, data)
            elif filename.endswith(".csv"):
                return self._write_csv_sheet(filename, data)
            else:
                return "Unsupported file format. Use .xlsx or .csv"
        except Exception as e:
            return f"Error writing sheet: {e}"

    def _write_excel_sheet(self, filename: str, data: List[List[str]]) -> str:
        """Write to Excel spreadsheet."""
        try:
            from openpyxl import Workbook  # type: ignore
        except ImportError:
            return "openpyxl not installed. Run: pip install openpyxl"

        wb = Workbook()  # type: ignore
        ws = wb.active
        if ws is None:
            return "Failed to create worksheet"

        for row in data:
            ws.append(row)

        wb.save(filename)
        return f"Data written to Excel sheet: {filename}"

    def _write_csv_sheet(self, filename: str, data: List[List[str]]) -> str:
        """Write to CSV file."""
        with open(filename, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerows(data)

        return f"Data written to CSV sheet: {filename}"

    # Presentations (PowerPoint)
    def create_presentation(
        self,
        filename: str,
        title: Optional[str] = None,
        content: Optional[str] = None,
    ) -> str:
        """Create a new PowerPoint presentation with an optional title slide."""
        if not filename:
            return "No filename provided."
        if not filename.lower().endswith(".pptx"):
            filename += ".pptx"
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return "python-pptx not installed. Run: pip install python-pptx"

        try:
            pres = Presentation()  # type: ignore
            # Title slide layout (0) if available
            layout = pres.slide_layouts[0]  # type: ignore
            slide = pres.slides.add_slide(layout)  # type: ignore
            if title:
                slide.shapes.title.text = title  # type: ignore
            if content:
                subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None  # type: ignore
                if subtitle is not None:
                    subtitle.text = content  # type: ignore
            pres.save(filename)  # type: ignore
            return f"Presentation created: {filename}"
        except Exception as e:
            return f"Error creating presentation: {e}"

    def add_presentation_slide(
        self,
        filename: str,
        title: Optional[str] = None,
        bullets: Optional[List[str]] = None,
    ) -> str:
        """Add a slide with title and bullet points to an existing presentation."""
        if not filename:
            return "No filename provided."
        if not os.path.exists(filename):
            return f"Presentation not found: {filename}"
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return "python-pptx not installed. Run: pip install python-pptx"

        try:
            pres = Presentation(filename)  # type: ignore
            layout = pres.slide_layouts[1]  # Title and Content
            slide = pres.slides.add_slide(layout)  # type: ignore
            if title:
                slide.shapes.title.text = title  # type: ignore
            if bullets:
                body = slide.shapes.placeholders[1].text_frame  # type: ignore
                body.clear()  # type: ignore
                first = True
                for b in bullets:
                    if first:
                        body.text = str(b)
                        first = False
                    else:
                        p = body.add_paragraph()
                        p.text = str(b)
                        p.level = 0
            pres.save(filename)  # type: ignore
            return f"Slide added to {filename}"
        except Exception as e:
            return f"Error adding slide: {e}"

    # Application Control
    def launch_app(self, app):
        if not app:
            return "No application specified."
        subprocess.Popen(app)
        return f"Launched application: {app}"

    def close_app(self, app):
        """Close application by name."""
        if not app:
            return "No application name provided."

        try:
            if os.name == "nt":  # Windows
                subprocess.run(
                    ["taskkill", "/F", "/IM", f"{app}.exe"],
                    capture_output=True,
                    text=True,
                )
                return f"Attempting to close {app}"
            else:  # Unix-like systems
                subprocess.run(
                    ["pkill", "-f", app], capture_output=True, text=True
                )
                return f"Attempting to close {app}"
        except Exception as e:
            return f"Error closing app: {e}"

    # Clipboard
    def clipboard_action(self, params):
        """Handle clipboard operations."""
        action = params.get("action", "")

        if action == "copy":
            return self.clipboard_copy(params.get("text", ""))
        elif action == "paste":
            return self.clipboard_paste()
        else:
            return "Unknown clipboard action. Use 'copy' or 'paste'."

    # Notifications
    def notify(self, message):
        """Send notification."""
        if not message:
            return "No message provided."

        return self.show_notification("M.I.A", message)

    # System Settings
    def system_setting(self, params):
        """Change system settings."""
        setting = params.get("setting", "")
        value = params.get("value", "")

        if not setting:
            return "No setting specified."

        try:
            if setting == "volume":
                return self._set_volume(value)
            elif setting == "brightness":
                return self._set_brightness(value)
            elif setting == "wifi":
                return self._control_wifi(value)
            else:
                return f"Setting '{setting}' not supported yet."
        except Exception as e:
            return f"Error changing setting: {e}"

    def _set_volume(self, value):
        """Set system volume."""
        try:
            if os.name == "nt":
                # Windows volume control
                subprocess.run(
                    ["nircmd", "setsysvolume", str(int(value) * 655.35)],
                    capture_output=True,
                )
                return f"Volume set to {value}%"
            else:
                # Linux volume control
                subprocess.run(
                    ["amixer", "set", "Master", f"{value}%"],
                    capture_output=True,
                )
                return f"Volume set to {value}%"
        except:
            return "Volume control not available on this system."

    def _set_brightness(self, value):
        """Set screen brightness."""
        try:
            if os.name == "nt":
                # Windows brightness control (requires powershell)
                cmd = f"(Get-WmiObject -Namespace root/WMI -Class WmiMonitorBrightnessMethods).WmiSetBrightness(1,{value})"
                subprocess.run(
                    ["powershell", "-Command", cmd], capture_output=True
                )
                return f"Brightness set to {value}%"
            else:
                # Linux brightness control
                subprocess.run(
                    ["xbacklight", "-set", str(value)], capture_output=True
                )
                return f"Brightness set to {value}%"
        except:
            return "Brightness control not available on this system."

    def _control_wifi(self, action):
        """Control WiFi connection."""
        try:
            if os.name == "nt":
                if action == "on":
                    subprocess.run(
                        [
                            "netsh",
                            "interface",
                            "set",
                            "interface",
                            "Wi-Fi",
                            "enabled",
                        ],
                        capture_output=True,
                    )
                    return "WiFi enabled"
                elif action == "off":
                    subprocess.run(
                        [
                            "netsh",
                            "interface",
                            "set",
                            "interface",
                            "Wi-Fi",
                            "disabled",
                        ],
                        capture_output=True,
                    )
                    return "WiFi disabled"
            else:
                if action == "on":
                    subprocess.run(
                        ["nmcli", "radio", "wifi", "on"], capture_output=True
                    )
                    return "WiFi enabled"
                elif action == "off":
                    subprocess.run(
                        ["nmcli", "radio", "wifi", "off"], capture_output=True
                    )
                    return "WiFi disabled"
        except:
            return "WiFi control not available on this system."

    def run_command(self, command):
        if not command:
            return "No command provided."
        result = subprocess.run(
            command, shell=True, capture_output=True, text=True
        )
        return result.stdout or result.stderr

    def run_sandboxed(self, params: Dict[str, Any]) -> str:
        sandbox = self._get_sandbox()
        if not sandbox:
            return "Sandbox provider not available. Install wasmtime/wasmer and enable the provider."

        module_path = params.get("module_path")
        wasm_bytes = params.get("wasm_bytes")
        wat_source = params.get("wat") or params.get("code")

        if wat_source and not wasm_bytes:
            try:
                from wasmtime import wat2wasm  # type: ignore
            except Exception:
                return "wat2wasm requires wasmtime. Install with 'pip install wasmtime'."
            try:
                wasm_bytes = wat2wasm(wat_source)
            except Exception as exc:  # pragma: no cover
                return f"Failed to compile WAT: {exc}"

        if not module_path and wasm_bytes is None:
            return "Provide 'module_path', 'wasm_bytes', or 'wat' to execute in the sandbox."

        stdin_data = params.get("stdin")
        if isinstance(stdin_data, bytes):
            stdin_bytes = stdin_data
        elif isinstance(stdin_data, str):
            stdin_bytes = stdin_data.encode("utf-8")
        else:
            stdin_bytes = None

        args_value = params.get("args") or []
        if isinstance(args_value, (str, bytes)):
            args_list = [str(args_value)]
        else:
            args_list = [str(arg) for arg in args_value]

        env_value = params.get("env") or {}
        if isinstance(env_value, dict):
            env_map = {str(k): str(v) for k, v in env_value.items()}
        else:
            env_map = {}

        try:
            result = sandbox.run(
                module_path=module_path,
                wasi_bytes=wasm_bytes,
                stdin=stdin_bytes,
                args=args_list,
                env=env_map,
            )
        except Exception as exc:  # pragma: no cover
            self.logger.error("Sandbox execution failed: %s", exc)
            return f"Sandbox execution failed: {exc}"

        return (
            f"Sandbox run {result.get('id')} exit_code={result.get('exit_code')} "
            f"engine={result.get('engine')} duration={result.get('duration_ms')}ms. "
            f"Logs: {result.get('log_file')}"
        )

    # Web Automation (Selenium)
    def web_automation(self, params: Dict[str, Any]) -> str:
        """Execute web automation tasks."""
        web_agent = self._get_web_agent()
        if not web_agent:
            return "Web agent not available. Install selenium and webdriver-manager."

        url = params.get("url")
        action = params.get("action", "browse")
        
        if not url and action == "browse":
            return "URL required for browsing."

        try:
            if action == "browse":
                # Simple browse/visit
                plan = [
                    {"action": "goto", "url": url},
                    {"action": "wait", "seconds": 2},
                    {"action": "screenshot", "name": "visit_result"}
                ]
                results = web_agent.run_plan(plan, headless=False)
                return f"Visited {url}. Results: {results}"
            
            elif action == "search":
                query = params.get("query")
                if not query:
                    return "Query required for search."
                
                # Construct a search plan (example for Google/DuckDuckGo)
                search_url = f"https://duckduckgo.com/?q={query.replace(' ', '+')}"
                plan = [
                    {"action": "goto", "url": search_url},
                    {"action": "wait", "seconds": 2},
                    {"action": "screenshot", "name": "search_result"}
                ]
                results = web_agent.run_plan(plan, headless=False)
                return f"Searched for '{query}'. Results: {results}"

            elif action == "execute_plan":
                plan = params.get("plan")
                if not plan:
                    return "Plan required for execution."
                results = web_agent.run_plan(plan, headless=False)
                return f"Executed web plan. Results: {results}"

            else:
                return f"Unknown web action: {action}"

        except Exception as e:
            return f"Web automation error: {e}"

    # Email (SMTP)
    def send_email(self, params):
        to = params.get("to")
        subject = params.get("subject", "(No Subject)")
        body = params.get("body", "")
        if not to:
            return "No recipient provided."

        # Resolve SMTP configuration from env/config with param overrides
        email_cfg = (self.config or {}).get("email", {})
        smtp_server = params.get("smtp_server") or email_cfg.get(
            "smtp_server", "smtp.gmail.com"
        )
        smtp_port = int(
            params.get("smtp_port") or email_cfg.get("smtp_port", 587)
        )
        from_addr = params.get("from") or email_cfg.get("username") or ""
        password = params.get("password") or email_cfg.get("password") or ""

        if not from_addr:
            return "Sender email not configured. Set EMAIL_USERNAME or provide 'from' param."

        msg = EmailMessage()
        msg["Subject"] = subject
        msg["From"] = from_addr
        msg["To"] = to
        msg.set_content(body)

        try:
            with smtplib.SMTP(smtp_server, smtp_port) as smtp:
                # Use STARTTLS if port typically supports it
                try:
                    smtp.starttls()
                except Exception:
                    pass
                if password:
                    smtp.login(from_addr, password)
                smtp.send_message(msg)
            return f"Email sent to {to}"
        except Exception as e:
            return f"Email error: {e}"

    # Calendar
    def calendar_event(self, params):
        """Create calendar event."""
        title = params.get("title", "")
        date = params.get("date", "")
        time = params.get("time", "")

        if not title:
            return "No event title provided."

        return "Google Calendar integration not configured."

    # Messaging
    def send_message(self, params):
        """Send message via various platforms."""
        platform = params.get("platform", "").lower()
        to = params.get("to", "")
        message = params.get("message", "")

        if not to or not message:
            return "Recipient and message required."

        try:
            if platform == "whatsapp":
                return self.send_whatsapp(params)
            elif platform == "telegram":
                return self._send_telegram(to, message)
            elif platform == "sms":
                return self._send_sms(to, message)
            else:
                return f"Platform '{platform}' not supported. Use 'whatsapp', 'telegram', or 'sms'."
        except Exception as e:
            return f"Error sending message: {e}"

    def _send_telegram(self, to, message):
        """Send Telegram message via configured messenger."""
        client = self._get_telegram_client()
        if not client:
            return "Telegram messenger not available. Configure credentials and install Telethon."
        try:
            return client.send_message(message, recipient=to)
        except Exception as exc:
            self.logger.error("Telegram send failed: %s", exc)
            return f"Error sending Telegram message: {exc}"

    def _send_sms(self, to, message):
        """Send SMS message."""
        return "SMS integration not configured."

    # Smart Home
    def smart_home(self, params):
        """Control smart home devices via Home Assistant."""
        device = params.get("device", "")
        action = params.get("action", "")

        if not device or not action:
            return "Device and action required."

        ha_url = self.config.get("smart_home", {}).get(
            "home_assistant_url", ""
        )
        ha_token = self.config.get("smart_home", {}).get(
            "home_assistant_token", ""
        )

        if not ha_url or not ha_token:
            return "Home Assistant URL and token not configured."

        try:
            headers = {"Authorization": f"Bearer {ha_token}"}
            url = f"{ha_url}/api/services/homeassistant/turn_{action}"
            data = {"entity_id": device}

            response = requests.post(url, headers=headers, json=data)
            if response.status_code == 200:
                return f"Smart home action: {action} on {device}"
            else:
                return f"Home Assistant error: {response.text}"
        except Exception as e:
            return f"Error controlling smart home: {e}"

    # Smart Home Integration
    def control_device(self, device_type: str, action: str, **kwargs) -> str:
        """Control smart home devices."""
        if not device_type or not action:
            return "Device type and action required."

        try:
            device_type = device_type.lower()
            action = action.lower()

            if device_type == "light":
                return self._control_light(action, **kwargs)
            elif device_type == "temperature":
                return self._control_temperature(action, **kwargs)
            elif device_type == "music":
                return self._control_music(action, **kwargs)
            elif device_type == "security":
                return self._control_security(action, **kwargs)
            else:
                return self._generic_device_control(
                    device_type, action, **kwargs
                )
        except Exception as e:
            return f"Error controlling device: {e}"

    def _control_light(self, action: str, **kwargs) -> str:
        """Control lighting system."""
        return "Smart home integration not configured."

    def _control_temperature(self, action: str, **kwargs) -> str:
        """Control temperature system."""
        return "Smart home integration not configured."

    def _control_music(self, action: str, **kwargs) -> str:
        """Control music system."""
        return "Smart home integration not configured."

    def _control_security(self, action: str, **kwargs) -> str:
        """Control security system."""
        return "Smart home integration not configured."

    def _generic_device_control(
        self, device_type: str, action: str, **kwargs
    ) -> str:
        """Handle generic device control."""
        return f"Controlling {device_type}: {action} with parameters {kwargs}"

    # System Integration
    def clipboard_copy(self, text: str) -> str:
        """Copy text to clipboard."""
        if not text:
            return "No text provided to copy."

        try:
            import pyperclip  # type: ignore
        except ImportError:
            return "pyperclip not installed. Run: pip install pyperclip"

        try:
            pyperclip.copy(text)  # type: ignore
            return f"Text copied to clipboard: {text[:50]}..."
        except Exception as e:
            return f"Error copying to clipboard: {e}"

    def clipboard_paste(self) -> str:
        """Get text from clipboard."""
        try:
            import pyperclip  # type: ignore
        except ImportError:
            return "pyperclip not installed. Run: pip install pyperclip"

        try:
            text = pyperclip.paste()  # type: ignore
            return f"Clipboard content: {text}"
        except Exception as e:
            return f"Error reading clipboard: {e}"

    def show_notification(self, title: str, message: str) -> str:
        """Show system notification."""
        if not title or not message:
            return "Title and message required for notification."

        try:
            from plyer import notification  # type: ignore
        except ImportError:
            return "plyer not installed. Run: pip install plyer"

        try:
            notification.notify(  # type: ignore
                title=title, message=message, app_name="M.I.A", timeout=10
            )
            return f"Notification sent: {title} - {message}"
        except Exception as e:
            return f"Error showing notification: {e}"

    def open_application(self, app_name: str) -> str:
        """Open an application."""
        if not app_name:
            return "No application name provided."

        try:
            if os.name == "nt":  # Windows
                os.startfile(app_name)
            elif os.name == "posix":  # Linux/Mac
                subprocess.run(
                    ["open", app_name]
                    if sys.platform == "darwin"
                    else ["xdg-open", app_name]
                )

            return f"Opening application: {app_name}"

        except Exception as e:
            return f"Error opening application: {e}"

    def get_system_info(self) -> str:
        """Get system information."""
        try:
            import platform

            info: Dict[str, Any] = {
                "system": platform.system(),
                "release": platform.release(),
                "version": platform.version(),
                "machine": platform.machine(),
                "processor": platform.processor(),
                "python_version": platform.python_version(),
            }

            try:
                import psutil  # type: ignore
            except ImportError:
                info["note"] = (
                    "psutil not installed. Run: pip install psutil for detailed system info"
                )
            else:
                info.update(
                    {
                        "cpu_count": psutil.cpu_count(),  # type: ignore
                        "memory": f"{psutil.virtual_memory().total / (1024**3):.2f} GB",  # type: ignore
                        "disk_usage": f"{psutil.disk_usage('/').percent}%",  # type: ignore
                    }
                )

            return json.dumps(info, indent=2)
        except Exception as e:
            return f"Error getting system info: {e}"

    def set_permission(self, action, allowed):
        self.permissions[action] = allowed

    def get_tool_descriptions(self) -> str:
        """Return a description of available tools for the LLM."""
        tools = [
            {
                "name": "web_search",
                "description": "Search the web for information.",
                "parameters": {"query": "The search query string"}
            },
            {
                "name": "web_scrape",
                "description": "Extract text content from a webpage.",
                "parameters": {"url": "The URL to scrape"}
            },
            {
                "name": "open_file",
                "description": "Read the contents of a file.",
                "parameters": {"path": "The file path"}
            },
            {
                "name": "create_file",
                "description": "Create a new file with content.",
                "parameters": {"path": "The file path", "content": "The content to write"}
            },
            {
                "name": "run_command",
                "description": "Run a shell command (use with caution).",
                "parameters": {"command": "The command to execute"}
            },
            {
                "name": "send_email",
                "description": "Send an email.",
                "parameters": {"to": "Recipient email", "subject": "Email subject", "body": "Email body"}
            },
            {
                "name": "desktop_open_app",
                "description": "Open a desktop application.",
                "parameters": {"app_name": "Name of the application"}
            },
            {
                "name": "store_memory",
                "description": "Store a memory for long-term retrieval.",
                "parameters": {"content": "The information to remember", "memory_type": "episodic or semantic"}
            },
            {
                "name": "search_memory",
                "description": "Search long-term memory.",
                "parameters": {"query": "The search query"}
            }
        ]
        
        description = "Available Tools:\n"
        for tool in tools:
            description += f"- {tool['name']}: {tool['description']}\n"
            description += f"  Parameters: {json.dumps(tool['parameters'])}\n"
        
        return description
